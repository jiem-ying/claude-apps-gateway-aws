#!/usr/bin/env python3
"""
Build the CONDENSED 2-slide "essence" governance deck (Anthropic-styled, 16:9).

A standalone ~3-minute companion to the full 14-slide deck (build_deck.py). Same
brand system, distilled to two slides:

  Slide 1  What it is · why govern · the AWS architecture
  Slide 2  What you can demo — Quota (spend caps) · Access (RBAC) · Visibility

The palette + helpers below are copied verbatim from build_deck.py so the two
decks look identical. This script does NOT touch the full deck or its output.

Run:  cd slides && python3 build_deck_2up.py
Out:  slides/claude-gateway-governance-2up.pptx

No network / LibreOffice needed — python-pptx writes .pptx natively.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ----------------------------------------------------------------- brand palette
# Anthropic-ish: ivory paper, near-black ink, clay/rust accent.
IVORY      = RGBColor(0xF0, 0xEE, 0xE6)   # paper background
INK        = RGBColor(0x14, 0x14, 0x13)   # near-black text
CLAY       = RGBColor(0xCC, 0x78, 0x5C)   # primary accent (rust/clay)
CLAY_DEEP  = RGBColor(0xA8, 0x54, 0x3A)   # darker clay for emphasis
SLATE      = RGBColor(0x3D, 0x3D, 0x3A)   # secondary ink
MUTE       = RGBColor(0x73, 0x70, 0x69)   # muted captions
SKY        = RGBColor(0x62, 0x8C, 0x9E)   # cool secondary (charts)
SAGE       = RGBColor(0x8A, 0x9A, 0x5B)   # green secondary (charts / "ok")
CARD       = RGBColor(0xFA, 0xF9, 0xF5)   # card fill (slightly lighter than paper)
CODE_BG    = RGBColor(0x22, 0x21, 0x1F)   # dark code panel
CODE_FG    = RGBColor(0xE8, 0xE6, 0xDE)   # code text
LINE       = RGBColor(0xD9, 0xD5, 0xC8)   # hairline dividers

SERIF = "Georgia"          # Tiempos stand-in (cross-platform)
SANS  = "Arial"            # Styrene stand-in (universal)

# 16:9
EMU_W, EMU_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width  = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]

# ----------------------------------------------------------------- helpers
def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = IVORY
    bg.line.fill.background()
    bg.shadow.inherit = False
    # send to back
    sp = bg._element; sp.getparent().remove(sp); s.shapes._spTree.insert(2, sp)
    return s

def _set_font(run, *, size, color, bold=False, italic=False, font=SANS):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    run.font.color.rgb = color

def textbox(s, x, y, w, h, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT, wrap=True):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = align
    return tb, tf

def add_para(tf, text, *, size, color, bold=False, italic=False, font=SANS,
             align=PP_ALIGN.LEFT, space_after=6, space_before=0, level=0,
             line_spacing=None, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.level = level
    if space_after is not None: p.space_after = Pt(space_after)
    if space_before is not None: p.space_before = Pt(space_before)
    if line_spacing: p.line_spacing = line_spacing
    r = p.add_run(); r.text = text
    _set_font(r, size=size, color=color, bold=bold, italic=italic, font=font)
    return p

def rect(s, x, y, w, h, fill, line=None, line_w=None, shadow=False, radius=False):
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = s.shapes.add_shape(shp_type, x, y, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = line_w or Pt(1)
    shp.shadow.inherit = False
    if radius:
        try:
            shp.adjustments[0] = 0.06
        except Exception:
            pass
    return shp

def hairline(s, x, y, w, color=LINE, weight=1.2):
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Pt(weight))
    ln.fill.solid(); ln.fill.fore_color.rgb = color
    ln.line.fill.background(); ln.shadow.inherit = False
    return ln

def kicker(s, text, x=Inches(0.9), y=Inches(0.62), color=CLAY):
    _, tf = textbox(s, x, y, Inches(11), Inches(0.4))
    add_para(tf, text.upper(), size=13, color=color, bold=True, font=SANS,
             first=True, space_after=0)

def title(s, text, x=Inches(0.9), y=Inches(1.0), w=Inches(11.5), size=34):
    _, tf = textbox(s, x, y, w, Inches(1.3))
    add_para(tf, text, size=size, color=INK, bold=False, font=SERIF,
             first=True, space_after=0, line_spacing=1.02)

def footer(s, idx):
    hairline(s, Inches(0.9), Inches(7.02), Inches(11.53), color=LINE, weight=1)
    _, tf = textbox(s, Inches(0.9), Inches(7.08), Inches(9), Inches(0.3))
    add_para(tf, "Claude Apps Gateway  ·  Governance", size=9, color=MUTE,
             font=SANS, first=True, space_after=0)
    _, tf2 = textbox(s, Inches(11.4), Inches(7.08), Inches(1.03), Inches(0.3),
                     align=PP_ALIGN.RIGHT)
    add_para(tf2, str(idx), size=9, color=MUTE, font=SANS, first=True,
             space_after=0, align=PP_ALIGN.RIGHT)

def code_panel(s, x, y, w, h, lines, title_text=None):
    """Dark code / terminal panel with monospace-ish text."""
    panel = rect(s, x, y, w, h, CODE_BG, radius=True)
    pad = Inches(0.22)
    ty = y + pad
    if title_text:
        _, tf = textbox(s, x + pad, y + Inches(0.12), w - 2*pad, Inches(0.3))
        add_para(tf, title_text, size=10, color=CLAY, bold=True, font=SANS,
                 first=True, space_after=0)
        ty = y + Inches(0.5)
    _, tf = textbox(s, x + pad, ty, w - 2*pad, h - (ty - y) - Inches(0.15))
    first = True
    for ln, col in lines:
        add_para(tf, ln if ln else " ", size=11.5, color=col, font="Courier New",
                 first=first, space_after=2, line_spacing=1.05)
        first = False
    return panel

def bullets(s, x, y, w, h, items, size=15, gap=10, marker_color=CLAY):
    """items: list of (bold_lead_or_None, text) ; renders a clay dash marker."""
    _, tf = textbox(s, x, y, w, h)
    first = True
    for it in items:
        lead, rest = it if isinstance(it, tuple) else (None, it)
        p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(gap)
        p.line_spacing = 1.08
        rm = p.add_run(); rm.text = "—  "
        _set_font(rm, size=size, color=marker_color, bold=True, font=SANS)
        if lead:
            rb = p.add_run(); rb.text = lead
            _set_font(rb, size=size, color=INK, bold=True, font=SANS)
            rt = p.add_run(); rt.text = rest
            _set_font(rt, size=size, color=SLATE, font=SANS)
        else:
            rt = p.add_run(); rt.text = rest
            _set_font(rt, size=size, color=SLATE, font=SANS)
    return tf

def flow_box(s, x, y, w, h, head, sub, fill, fg, border):
    """A single labelled node in the architecture flow strip."""
    shp = rect(s, x, y, w, h, fill, line=(None if border is None else border),
               line_w=Pt(1.5), radius=True)
    tf = shp.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(4); tf.margin_right = Pt(4)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = head
    _set_font(r, size=13.5, color=fg, bold=True, font=SERIF)
    if sub:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(3)
        r2 = p2.add_run(); r2.text = sub
        _set_font(r2, size=9, color=(IVORY if fill == CLAY else MUTE), font=SANS)
    return shp

def arrow(s, x, y_mid, w):
    """Clay connector arrow between two flow boxes, centred on y_mid."""
    rect(s, x + Inches(0.04), y_mid - Pt(1.5), w - Inches(0.22), Pt(3), CLAY_DEEP)
    tip = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                             x + w - Inches(0.20), y_mid - Inches(0.085),
                             Inches(0.17), Inches(0.17))
    tip.rotation = 90; tip.fill.solid(); tip.fill.fore_color.rgb = CLAY_DEEP
    tip.line.fill.background(); tip.shadow.inherit = False

# ============================================================ SLIDE 1 — what · why · architecture
s = slide()
# clay side band
rect(s, 0, 0, Inches(0.28), EMU_H, CLAY)
kicker(s, "Governing Claude Code", x=Inches(0.95))
title(s, "One credential, in your account, under your policy", x=Inches(0.95), size=32)

_, tf = textbox(s, Inches(0.95), Inches(1.78), Inches(11.6), Inches(0.9))
add_para(tf, "The Claude Apps Gateway is a self-hosted control plane on AWS. Developers sign in with "
             "corporate SSO; the gateway routes every request to Amazon Bedrock on one IAM role it holds — "
             "so there are no AWS keys on any laptop, and offboarding is just removing someone from your IdP.",
         size=13.5, color=SLATE, font=SANS, first=True, space_after=0, line_spacing=1.2)

# --- architecture flow strip ---
_, tf = textbox(s, Inches(0.95), Inches(2.9), Inches(11), Inches(0.3))
add_para(tf, "THE PATH EVERY REQUEST TAKES", size=11, color=MUTE, bold=True, font=SANS,
         first=True, space_after=0)

nodes = [
    ("Developer\nClaude Code", "Corporate SSO", IVORY, INK, CLAY),
    ("OIDC IdP", "Cognito or BYO\n(Okta · Entra)", IVORY, INK, CLAY),
    ("Internal ALB", "Private · IPv4", IVORY, INK, CLAY),
    ("Gateway\nECS Fargate", "Enforces policy", CLAY, IVORY, None),
    ("Amazon\nBedrock", "Claude, your region", IVORY, INK, CLAY),
]
bx = Inches(0.95); bw = Inches(2.06); bh = Inches(1.5); by = Inches(3.28)
gapx = Inches(0.42)
for i, (head, sub, fill, fg, border) in enumerate(nodes):
    x = bx + i * (bw + gapx)
    flow_box(s, x, by, bw, bh, head, sub, fill, fg, border)
    if i < len(nodes) - 1:
        arrow(s, x + bw, by + bh/2, gapx)

_, tf = textbox(s, Inches(0.95), Inches(5.05), Inches(11.6), Inches(0.55))
add_para(tf, "Everything beneath the token stays inside your VPC — ", size=12.5, color=INK, bold=True,
         font=SANS, first=True, space_after=0, line_spacing=1.15)
add_para(tf, "RDS PostgreSQL holds auth & spend state, an optional collector forwards metrics over native "
             "OTLP to CloudWatch, all in one account. Unless you deliberately add it, traffic never leaves your environment.",
         size=12.5, color=SLATE, font=SANS, space_before=2, space_after=0, line_spacing=1.18)

# --- why govern: three-risk strip ---
_, tf = textbox(s, Inches(0.95), Inches(5.95), Inches(11), Inches(0.3))
add_para(tf, "WHY GOVERN — ONE SHARED CREDENTIAL CONCENTRATES THREE RISKS", size=11, color=MUTE,
         bold=True, font=SANS, first=True, space_after=0)

risks = [
    ("Runaway spend", "one looping agent can burn the whole commitment", CLAY),
    ("Access sprawl", "every dev can reach every model and every tool", SKY),
    ("No attribution", "the invoice arrives as one anonymous lump sum", SAGE),
]
rx = Inches(0.95); rw = Inches(3.71); rgap = Inches(0.2); ry = Inches(6.28); rh = Inches(0.62)
for i, (h, body, col) in enumerate(risks):
    x = rx + i * (rw + rgap)
    rect(s, x, ry, rw, rh, CARD, line=LINE, line_w=Pt(1), radius=True)
    rect(s, x, ry, Inches(0.09), rh, col, radius=False)
    _, tf = textbox(s, x + Inches(0.24), ry, rw - Inches(0.4), rh, anchor=MSO_ANCHOR.MIDDLE)
    add_para(tf, h, size=13, color=INK, bold=True, font=SERIF, first=True, space_after=1)
    add_para(tf, body, size=10.5, color=SLATE, font=SANS, space_after=0, line_spacing=1.05)
footer(s, 1)

# ============================================================ SLIDE 2 — what you can demo
s = slide()
kicker(s, "What you can demo")
title(s, "Configure it, cap it, and see who did what")

_, tf = textbox(s, Inches(0.9), Inches(1.72), Inches(11.5), Inches(0.55))
add_para(tf, "Three governance controls, all enforced server-side at the API — not merely hidden in the client.",
         size=14, color=SLATE, font=SANS, first=True, space_after=0, line_spacing=1.15)

cards = [
    ("01", "Quota", "Real-time spend caps", CLAY, [
        ("Per-user, per-group, or org daily / weekly / monthly USD budgets, checked live on every request."),
        ("DEMO: over-cap request returns ", "429 billing_error"),
        ("Admin raises it via the spend_limits API — next request succeeds. No redeploy."),
        ("Estimated at list price: a circuit breaker, not an invoice."),
    ]),
    ("02", "Access", "OIDC-group RBAC", SKY, [
        ("An identity group (Cognito, Okta, Entra) maps to its allowed models and tools."),
        ("", "ENFORCE_MODELS", " org allowlist → off-list model rejected with 400."),
        ("", "DENY_TOOL_GROUP", " for a group withholds a built-in like WebFetch."),
        ("Gateway gates access; it does not push MCP servers."),
    ]),
    ("03", "Visibility", "Attribution & alerts", SAGE, [
        ("Every metric and audit event carries the developer's identity."),
        ("Managed Coding Agent Insights dashboard: cost by user, team, model."),
        ("Governance forensics via Logs Insights over /aws/claude-gateway/events."),
        ("PromQL alarms on daily cost and per-user cost."),
    ]),
]
cx = Inches(0.9); cw = Inches(3.71); gap = Inches(0.2); cy = Inches(2.5); ch = Inches(3.95)
for i, (num, name, tag, col, lines) in enumerate(cards):
    x = cx + i * (cw + gap)
    rect(s, x, cy, cw, ch, CARD, line=LINE, line_w=Pt(1), radius=True)
    rect(s, x, cy, cw, Inches(0.14), col, radius=False)
    _, tf = textbox(s, x + Inches(0.28), cy + Inches(0.36), cw - Inches(0.56), ch - Inches(0.6))
    add_para(tf, num, size=24, color=col, bold=True, font=SERIF, first=True, space_after=2)
    add_para(tf, name, size=20, color=INK, bold=True, font=SERIF, space_after=1)
    add_para(tf, tag, size=12.5, color=col, bold=True, font=SANS, space_after=12)
    for item in lines:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT; p.space_after = Pt(8); p.line_spacing = 1.1
        rm = p.add_run(); rm.text = "—  "
        _set_font(rm, size=11.5, color=col, bold=True, font=SANS)
        if isinstance(item, tuple):
            # (plain, mono) or (plain, mono, plain-tail) — mono chunk in Courier
            rt = p.add_run(); rt.text = item[0]
            _set_font(rt, size=11.5, color=SLATE, font=SANS)
            rc = p.add_run(); rc.text = item[1]
            _set_font(rc, size=11, color=INK, bold=True, font="Courier New")
            if len(item) > 2:
                rtail = p.add_run(); rtail.text = item[2]
                _set_font(rtail, size=11.5, color=SLATE, font=SANS)
        else:
            rt = p.add_run(); rt.text = item
            _set_font(rt, size=11.5, color=SLATE, font=SANS)

_, tf = textbox(s, Inches(0.9), Inches(6.62), Inches(11.5), Inches(0.4))
add_para(tf, "One self-hosted gateway on AWS — running on your credential, under your policy, in full view.",
         size=14, color=INK, italic=True, font=SERIF, first=True, space_after=0, align=PP_ALIGN.CENTER)
footer(s, 2)

# ----------------------------------------------------------------- save
import os
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "claude-gateway-governance-2up.pptx")
prs.save(out)
print(f"wrote {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
