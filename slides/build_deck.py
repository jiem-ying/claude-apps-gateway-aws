#!/usr/bin/env python3
"""
Build the Claude Apps Gateway governance demo deck (Anthropic-styled, 16:9).

Covers three high-impact governance features:
  1. Quota / spend management  (configure -> track -> cap, with hard 429 enforcement)
  2. Group RBAC                (per-IdP-group model + tool access, enforced server-side)
  3. Observability / audit     (identity-stamped cost + governance dashboard + alarms)

Audience: mixed exec + technical. Each pillar leads with business value, then shows
the real config / API call.

Run:  cd slides && python3 build_deck.py
Out:  slides/claude-gateway-governance.pptx

No network / LibreOffice needed — python-pptx writes .pptx natively.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ----------------------------------------------------------------- brand palette
# Anthropic-ish: ivory paper, near-black ink, clay/rust accent.
IVORY      = RGBColor(0xF0, 0xEE, 0xE6)   # paper background
INK        = RGBColor(0x14, 0x14, 0x13)   # near-black text
CLAY       = RGBColor(0xCC, 0x78, 0x5C)   # primary accent (rust/clay)
CLAY_DEEP  = RGBColor(0xA8, 0x54, 0x3A)   # darker clay for emphasis
SLATE      = RGBColor(0x3D, 0x3D, 0x3A)   # secondary ink
MUTE       = RGBColor(0x73, 0x70, 0x69)   # muted captions
SKY        = RGBColor(0x62, 0x8C, 0x9E)   # cool secondary (charts)
SAGE       = RGBColor(0x8A, 0x9A, 0x5B)   # green secondary (charts / "ok")
CARD       = RGBColor(0xFA, 0xF9, 0xF5)   # card fill (slightly lighter than paper)
CODE_BG    = RGBColor(0x22, 0x21, 0x1F)   # dark code panel
CODE_FG    = RGBColor(0xE8, 0xE6, 0xDE)   # code text
LINE       = RGBColor(0xD9, 0xD5, 0xC8)   # hairline dividers

SERIF = "Georgia"          # Tiempos stand-in (cross-platform)
SANS  = "Arial"            # Styrene stand-in (universal)

# 16:9
EMU_W, EMU_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width  = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]

# ----------------------------------------------------------------- helpers
def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = IVORY
    bg.line.fill.background()
    bg.shadow.inherit = False
    # send to back
    sp = bg._element; sp.getparent().remove(sp); s.shapes._spTree.insert(2, sp)
    return s

def _set_font(run, *, size, color, bold=False, italic=False, font=SANS):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    run.font.color.rgb = color

def textbox(s, x, y, w, h, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT, wrap=True):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = align
    return tb, tf

def add_para(tf, text, *, size, color, bold=False, italic=False, font=SANS,
             align=PP_ALIGN.LEFT, space_after=6, space_before=0, level=0,
             line_spacing=None, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.level = level
    if space_after is not None: p.space_after = Pt(space_after)
    if space_before is not None: p.space_before = Pt(space_before)
    if line_spacing: p.line_spacing = line_spacing
    r = p.add_run(); r.text = text
    _set_font(r, size=size, color=color, bold=bold, italic=italic, font=font)
    return p

def rect(s, x, y, w, h, fill, line=None, line_w=None, shadow=False, radius=False):
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = s.shapes.add_shape(shp_type, x, y, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = line_w or Pt(1)
    shp.shadow.inherit = False
    if radius:
        try:
            shp.adjustments[0] = 0.06
        except Exception:
            pass
    return shp

def hairline(s, x, y, w, color=LINE, weight=1.2):
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Pt(weight))
    ln.fill.solid(); ln.fill.fore_color.rgb = color
    ln.line.fill.background(); ln.shadow.inherit = False
    return ln

def kicker(s, text, x=Inches(0.9), y=Inches(0.62), color=CLAY):
    _, tf = textbox(s, x, y, Inches(11), Inches(0.4))
    add_para(tf, text.upper(), size=13, color=color, bold=True, font=SANS,
             first=True, space_after=0)

def title(s, text, x=Inches(0.9), y=Inches(1.0), w=Inches(11.5), size=34):
    _, tf = textbox(s, x, y, w, Inches(1.3))
    add_para(tf, text, size=size, color=INK, bold=False, font=SERIF,
             first=True, space_after=0, line_spacing=1.02)

def footer(s, idx):
    hairline(s, Inches(0.9), Inches(7.02), Inches(11.53), color=LINE, weight=1)
    _, tf = textbox(s, Inches(0.9), Inches(7.08), Inches(9), Inches(0.3))
    add_para(tf, "Claude Apps Gateway  ·  Governance", size=9, color=MUTE,
             font=SANS, first=True, space_after=0)
    _, tf2 = textbox(s, Inches(11.4), Inches(7.08), Inches(1.03), Inches(0.3),
                     align=PP_ALIGN.RIGHT)
    add_para(tf2, str(idx), size=9, color=MUTE, font=SANS, first=True,
             space_after=0, align=PP_ALIGN.RIGHT)

def code_panel(s, x, y, w, h, lines, title_text=None):
    """Dark code / terminal panel with monospace-ish text."""
    panel = rect(s, x, y, w, h, CODE_BG, radius=True)
    pad = Inches(0.22)
    ty = y + pad
    if title_text:
        _, tf = textbox(s, x + pad, y + Inches(0.12), w - 2*pad, Inches(0.3))
        add_para(tf, title_text, size=10, color=CLAY, bold=True, font=SANS,
                 first=True, space_after=0)
        ty = y + Inches(0.5)
    _, tf = textbox(s, x + pad, ty, w - 2*pad, h - (ty - y) - Inches(0.15))
    first = True
    for ln, col in lines:
        add_para(tf, ln if ln else " ", size=11.5, color=col, font="Courier New",
                 first=first, space_after=2, line_spacing=1.05)
        first = False
    return panel

def bullets(s, x, y, w, h, items, size=15, gap=10, marker_color=CLAY):
    """items: list of (bold_lead_or_None, text) ; renders a clay dash marker."""
    _, tf = textbox(s, x, y, w, h)
    first = True
    for it in items:
        lead, rest = it if isinstance(it, tuple) else (None, it)
        p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(gap)
        p.line_spacing = 1.08
        rm = p.add_run(); rm.text = "—  "
        _set_font(rm, size=size, color=marker_color, bold=True, font=SANS)
        if lead:
            rb = p.add_run(); rb.text = lead
            _set_font(rb, size=size, color=INK, bold=True, font=SANS)
            rt = p.add_run(); rt.text = rest
            _set_font(rt, size=size, color=SLATE, font=SANS)
        else:
            rt = p.add_run(); rt.text = rest
            _set_font(rt, size=size, color=SLATE, font=SANS)
    return tf

# ============================================================ SLIDE 1 — title
s = slide()
# clay side band
rect(s, 0, 0, Inches(0.28), EMU_H, CLAY)
_, tf = textbox(s, Inches(0.95), Inches(2.35), Inches(11.5), Inches(2.4))
add_para(tf, "GOVERNING CLAUDE CODE", size=15, color=CLAY, bold=True, font=SANS,
         first=True, space_after=10)
add_para(tf, "Governance at enterprise scale", size=46, color=INK, font=SERIF,
         space_after=6, line_spacing=1.0)
add_para(tf, "Quota control · Access control · Full visibility", size=20,
         color=SLATE, font=SERIF, italic=True, space_before=4)
_, tf2 = textbox(s, Inches(0.95), Inches(6.2), Inches(11.3), Inches(0.7))
add_para(tf2, "A self-hosted gateway on AWS that routes Claude Code through Amazon Bedrock, using your own "
              "identity provider, network, and telemetry.",
         size=13.5, color=MUTE, font=SANS, first=True, space_after=0, line_spacing=1.2)

# ============================================================ SLIDE 2 — the problem
s = slide()
kicker(s, "Why governance")
title(s, "One shared credential. Hundreds of agents.")
_, tf = textbox(s, Inches(0.9), Inches(1.85), Inches(11.5), Inches(0.7))
add_para(tf, "The gateway routes every developer's inference through one upstream credential. "
             "That's the point — no API keys on laptops — but it concentrates three risks:",
         size=15, color=SLATE, font=SANS, first=True, space_after=0, line_spacing=1.1)

cards = [
    ("Runaway spend", "Because every request bills to the same credential, a single looping "
        "agent can quietly consume the organization's entire commitment before anyone notices.", CLAY),
    ("Access sprawl", "Without policy, every developer can reach every model and every tool, "
        "so a contractor and a senior engineer look exactly the same to the upstream.", SKY),
    ("No attribution", "The provider invoice arrives as one lump sum, leaving you unable to say "
        "which team or which person actually drove the cost.", SAGE),
]
cx = Inches(0.9); cw = Inches(3.71); gap = Inches(0.2); cy = Inches(2.95); ch = Inches(3.3)
for i, (h, body, col) in enumerate(cards):
    x = cx + i * (cw + gap)
    rect(s, x, cy, cw, ch, CARD, line=LINE, line_w=Pt(1), radius=True)
    rect(s, x, cy, cw, Inches(0.14), col, radius=False)
    _, tf = textbox(s, x + Inches(0.28), cy + Inches(0.42), cw - Inches(0.56), ch - Inches(0.7))
    add_para(tf, h, size=19, color=INK, bold=True, font=SERIF, first=True, space_after=10)
    add_para(tf, body, size=13.5, color=SLATE, font=SANS, space_after=0, line_spacing=1.14)
footer(s, 2)

# ============================================================ SLIDE 3 — architecture
s = slide()
kicker(s, "How it fits")
title(s, "The gateway sits between developers and Bedrock")
# simple flow diagram
boxes = [
    ("Developer\nClaude Code", "Signs in with SSO", IVORY, INK),
    ("Claude Apps\nGateway", "Holds the credential and enforces policy", CLAY, IVORY),
    ("Amazon\nBedrock", "Serves Claude in your region", IVORY, INK),
]
bx = Inches(1.15); bw = Inches(3.0); bh = Inches(1.7); by = Inches(2.75); gapx = Inches(1.55)
centers = []
for i, (h, sub, fill, fg) in enumerate(boxes):
    x = bx + i * (bw + gapx)
    shp = rect(s, x, by, bw, bh, fill, line=(None if fill==CLAY else CLAY), line_w=Pt(1.5), radius=True)
    tf = shp.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = h
    _set_font(r, size=17, color=fg, bold=True, font=SERIF)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(6)
    r2 = p2.add_run(); r2.text = sub
    _set_font(r2, size=10.5, color=(IVORY if fill==CLAY else MUTE), font=SANS)
    centers.append((x + bw, by + bh/2))
    if i < len(boxes) - 1:
        ax = x + bw; aw = gapx
        arr = rect(s, ax + Inches(0.15), by + bh/2 - Pt(1.5), aw - Inches(0.3), Pt(3), CLAY_DEEP)
        tip = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, ax + aw - Inches(0.2),
                                 by + bh/2 - Inches(0.09), Inches(0.18), Inches(0.18), )
        tip.rotation = 90; tip.fill.solid(); tip.fill.fore_color.rgb = CLAY_DEEP
        tip.line.fill.background(); tip.shadow.inherit = False

_, tf = textbox(s, Inches(1.15), Inches(4.95), Inches(11), Inches(1.3))
add_para(tf, "Everything beneath the token stays inside your account.", size=15, color=INK, bold=True,
         font=SANS, first=True, space_after=8)
add_para(tf, "The gateway runs behind a private internal load balancer, signs users in against your own "
             "identity provider, and stores its state in your database. Unless you deliberately add the "
             "Anthropic API as an upstream, none of that traffic ever leaves your environment.",
         size=13.5, color=SLATE, font=SANS, space_after=0, line_spacing=1.18)
footer(s, 3)

# ============================================================ SLIDE 4 — three pillars
s = slide()
kicker(s, "The governance model")
title(s, "Three pillars, enforced server-side")
pillars = [
    ("01", "Quota", "Cap spend before it happens",
     "You give each person, team, or the whole organization a budget, and the gateway refuses any "
     "request that would exceed it — a circuit breaker on a shared bill.", CLAY),
    ("02", "Access", "The right models and tools per team",
     "Each identity group maps to its own set of allowed models and tools, and the gateway enforces "
     "that choice at the API rather than merely hiding it in the interface.", SKY),
    ("03", "Visibility", "Know who spent what, and what they did",
     "Every metric and audit event carries the developer's identity, so you can attribute cost, review "
     "activity, and be alerted the moment spending looks unusual.", SAGE),
]
cx = Inches(0.9); cw = Inches(3.71); gap = Inches(0.2); cy = Inches(2.1); ch = Inches(4.35)
for i, (num, h, tag, body, col) in enumerate(pillars):
    x = cx + i * (cw + gap)
    rect(s, x, cy, cw, ch, CARD, line=LINE, line_w=Pt(1), radius=True)
    _, tf = textbox(s, x + Inches(0.3), cy + Inches(0.35), cw - Inches(0.6), ch - Inches(0.6))
    add_para(tf, num, size=30, color=col, bold=True, font=SERIF, first=True, space_after=4)
    add_para(tf, h, size=22, color=INK, bold=True, font=SERIF, space_after=2)
    add_para(tf, tag, size=13.5, color=col, bold=True, font=SANS, space_after=12)
    add_para(tf, body, size=13.5, color=SLATE, font=SANS, space_after=0, line_spacing=1.16)
footer(s, 4)

# ============================================================ SLIDE 5 — Quota value
s = slide()
rect(s, 0, 0, Inches(0.28), EMU_H, CLAY)
kicker(s, "Pillar 01 · Quota management")
title(s, "Cap runaway spend on a credential everyone shares")
bullets(s, Inches(0.9), Inches(2.2), Inches(6.9), Inches(4.2), [
    ("You set the ceiling. ", "Each developer, each team, or the whole organization is given a daily, "
        "weekly, or monthly budget that reflects how you actually want the shared credential used."),
    ("The gateway enforces it live. ", "It checks the running total on every request, so a single "
        "runaway agent is stopped in the moment rather than discovered on next month's invoice."),
    ("Raising a limit takes seconds. ", "When someone legitimately needs more room, an administrator "
        "lifts their cap with one API call and the very next request succeeds — no redeployment required."),
], size=15.5, gap=18)
# stat card
rect(s, Inches(8.15), Inches(2.3), Inches(4.25), Inches(3.5), INK, radius=True)
_, tf = textbox(s, Inches(8.5), Inches(2.75), Inches(3.6), Inches(3.0))
add_para(tf, "429", size=64, color=CLAY, bold=True, font=SERIF, first=True, space_after=0, align=PP_ALIGN.CENTER)
add_para(tf, "billing_error", size=15, color=CODE_FG, font="Courier New", space_after=14, align=PP_ALIGN.CENTER)
add_para(tf, "This is what a developer receives the moment they cross their cap, and it clears "
             "automatically once the period resets or an administrator raises the limit.",
         size=13, color=IVORY, font=SANS, space_after=0, align=PP_ALIGN.CENTER, line_spacing=1.2)
footer(s, 5)

# ============================================================ SLIDE 6 — Quota how / API
s = slide()
kicker(s, "Pillar 01 · Quota management")
title(s, "Configure once, set caps with a simple API")
bullets(s, Inches(0.9), Inches(2.05), Inches(5.5), Inches(4.6), [
    ("You enable it once at deploy time. ", "Setting a single flag turns on the admin API, and the "
        "secret write key is injected securely from Secrets Manager."),
    ("You then set caps by scope. ", "The same call targets one developer, an entire identity group, "
        "or the organization as a whole, with the amount given in US cents."),
    ("Availability stays in your control. ", "If the database is briefly unreachable, the gateway keeps "
        "serving by default; you can instead choose to fail closed so no spend ever goes unmetered."),
], size=14.5, gap=16)

code_panel(s, Inches(6.65), Inches(1.95), Inches(5.75), Inches(4.65), [
    ("# 1. Enable enforcement at deploy time", MUTE),
    ("export ENABLE_SPEND_CAPS=true", CODE_FG),
    ("export GATEWAY_ADMIN_WRITE_KEY_ARN=arn:...", CODE_FG),
    ("./deploy.sh", CODE_FG),
    ("", CODE_FG),
    ("# 2. $500 / developer / month, org-wide", MUTE),
    ("curl -X POST $GW/v1/organizations/\\", SAGE),
    ("  spend_limits \\", SAGE),
    ('  -H "x-api-key: $ADMIN_KEY" \\', CODE_FG),
    ("  -d '{\"scope\":{\"type\":", CODE_FG),
    ("       \"organization\"},", CODE_FG),
    ('       "amount":"50000",', CODE_FG),
    ('       "period":"monthly"}\'', CODE_FG),
    ("", CODE_FG),
    ("# amount is USD cents · null = unlimited", MUTE),
], title_text="terminal")
footer(s, 6)

# ============================================================ SLIDE 7 — Quota resolution + effective
s = slide()
kicker(s, "Pillar 01 · Quota management")
title(s, "Which budget applies to a given developer")
_, tf = textbox(s, Inches(0.9), Inches(1.85), Inches(6.7), Inches(1.5))
add_para(tf, "A developer can be covered by several budgets at once — one they were given personally, "
             "one attached to each team they belong to, and the organization-wide default. Rather than "
             "adding these together, the gateway picks the single budget that should govern them, working "
             "from the most specific rule to the most general.",
         size=14, color=SLATE, font=SANS, first=True, space_after=0, line_spacing=1.22)

steps = [
    ("If they have a personal budget, that wins.",
     "A limit set on an individual always takes precedence, which is exactly how you exempt or specially "
     "provision one person."),
    ("Otherwise, their tightest team budget applies.",
     "When someone belongs to several teams, the gateway defaults to the most restrictive of those team "
     "budgets, and you can flip this to the most generous if you prefer."),
    ("Failing that, the organization default applies.",
     "Everyone who matched nothing more specific inherits the company-wide ceiling."),
    ("And if none of those exist, they are unlimited.",
     "No budget at any level simply means no cap is enforced for that person."),
]
y = Inches(3.35)
for h, sub in steps:
    _, tf = textbox(s, Inches(0.9), y, Inches(6.75), Inches(0.85))
    add_para(tf, h, size=13.5, color=INK, bold=True, font=SANS, first=True, space_after=2)
    add_para(tf, sub, size=11.5, color=SLATE, font=SANS, space_after=0, line_spacing=1.12)
    y += Inches(0.9)

# worked-example panel
code_panel(s, Inches(7.95), Inches(2.05), Inches(4.45), Inches(4.55), [
    ("A worked example", CLAY),
    ("", CODE_FG),
    ("The organization sets a", CODE_FG),
    ("$20 daily default, and the", CODE_FG),
    ("contractors team is given a", CODE_FG),
    ("tighter $5 daily budget.", CODE_FG),
    ("", CODE_FG),
    ("For that period the more", CODE_FG),
    ("restrictive team budget is", CODE_FG),
    ("the one that governs, so a", CODE_FG),
    ("contractor is held to $5 a", CODE_FG),
    ("day, not the $20 default.", MUTE),
], title_text=None)
footer(s, 7)

# ============================================================ SLIDE 8 — RBAC value
s = slide()
rect(s, 0, 0, Inches(0.28), EMU_H, SKY)
kicker(s, "Pillar 02 · Access control", color=SKY)
title(s, "The right models and tools for each team")
bullets(s, Inches(0.9), Inches(2.2), Inches(6.9), Inches(4.2), [
    ("Policy follows the team, not the person. ", "You attach rules to the groups your identity "
        "provider already manages, so a new contractor picks up the contractor policy from their group "
        "membership at sign-in — there is no per-person setup to maintain."),
    ("The limit is enforced, not merely suggested. ", "When a policy withholds a model, the gateway "
        "rejects any request for it outright, so a modified or scripted client cannot reach around the "
        "interface to use it anyway."),
    ("Tools are governed the same way. ", "You can withhold a risky tool from one group while leaving "
        "their choice of models completely intact, matching capability to trust."),
], size=15.5, gap=17)
# mini matrix
rect(s, Inches(8.15), Inches(2.3), Inches(4.25), Inches(3.7), CARD, line=LINE, line_w=Pt(1), radius=True)
_, tf = textbox(s, Inches(8.45), Inches(2.55), Inches(3.7), Inches(3.3))
add_para(tf, "One policy per team", size=14, color=INK, bold=True, font=SERIF, first=True, space_after=12)
rows = [("Core engineering", "Every model, every tool", SAGE),
        ("Data science", "Opus and Sonnet, no shell access", SKY),
        ("Contractors", "Sonnet only, web fetch withheld", CLAY)]
for name, desc, col in rows:
    add_para(tf, name, size=13, color=col, bold=True, font=SANS, space_after=1, space_before=8)
    add_para(tf, desc, size=11.5, color=SLATE, font=SANS, space_after=0, line_spacing=1.1)
footer(s, 8)

# ============================================================ SLIDE 9 — RBAC how / YAML
s = slide()
kicker(s, "Pillar 02 · Access control", color=SKY)
title(s, "Policies are plain config, first match wins")
code_panel(s, Inches(0.9), Inches(1.95), Inches(6.5), Inches(4.65), [
    ("managed:", CODE_FG),
    ("  policies:", CODE_FG),
    ("    # specific groups first", MUTE),
    ("    - match: {groups: [contractors]}", SKY),
    ("      cli:", CODE_FG),
    ("        availableModels:", CODE_FG),
    ("          [claude-sonnet-4-6]", CODE_FG),
    ("        permissions:", CODE_FG),
    ('          deny: ["WebFetch"]', CLAY),
    ("    # catch-all base, last", MUTE),
    ("    - match: {}", SAGE),
    ("      cli:", CODE_FG),
    ("        availableModels:", CODE_FG),
    ("          [claude-opus-4-8,", CODE_FG),
    ("           claude-sonnet-4-6,", CODE_FG),
    ("           claude-haiku-4-5]", CODE_FG),
], title_text="gateway.yaml")
bullets(s, Inches(7.7), Inches(2.05), Inches(4.75), Inches(4.5), [
    ("The empty match is your baseline. ", "It applies to everyone, and each group's policy only needs "
        "to state what it changes from that baseline."),
    ("Org-wide restrictions cannot be lost. ", "A denial or audit rule set for the whole organization is "
        "always combined with a group's own rules, so a team policy can add restrictions but never quietly "
        "drop one the organization requires."),
    ("Changes reach clients on their own. ", "After you update a policy and redeploy, signed-in clients "
        "adopt it at their next hourly check, and a change of team membership takes effect the next time "
        "the developer signs in."),
], size=13.5, gap=15, marker_color=SKY)
footer(s, 9)

# ============================================================ SLIDE 10 — Observability value
s = slide()
rect(s, 0, 0, Inches(0.28), EMU_H, SAGE)
kicker(s, "Pillar 03 · Visibility", color=SAGE)
title(s, "Know who spent what — and what they did")
bullets(s, Inches(0.9), Inches(2.2), Inches(6.7), Inches(4.2), [
    ("Every number carries a name. ", "The developer's identity travels with each metric as an attribute, so "
        "cost and activity attribute to a person and a team without anyone configuring anything."),
    ("The usage dashboard is managed for you. ", "Metrics flow to CloudWatch and auto-populate the built-in "
        "Coding Agent Insights dashboard — cost and tokens by user, team, and model, plus adoption — with "
        "nothing to build. The gateway's own dashboard is then just the governance and audit companion."),
    ("Alerts can trigger action, not just email. ", "A threshold or anomaly raises a notification that "
        "can reach a person on Slack or page an on-call engineer, and can equally invoke a function that "
        "adjusts a budget on its own."),
], size=14.5, gap=17)

# mini bar chart: cost by team
chart_x = Inches(8.0); chart_y = Inches(2.35); chart_w = Inches(4.4); chart_h = Inches(3.55)
rect(s, chart_x, chart_y, chart_w, chart_h, CARD, line=LINE, line_w=Pt(1), radius=True)
_, tf = textbox(s, chart_x + Inches(0.3), chart_y + Inches(0.22), chart_w - Inches(0.6), Inches(0.4))
add_para(tf, "Coding Agent Insights — cost by team", size=13, color=INK, bold=True, font=SERIF, first=True, space_after=0)
teams = [("Core eng", 0.92, CLAY), ("Data sci", 0.61, SKY), ("Platform", 0.44, SAGE),
         ("Contractors", 0.22, MUTE)]
base_x = chart_x + Inches(1.5)
bar_area = Inches(2.35)
bar_h = Inches(0.42); by0 = chart_y + Inches(0.85); step = Inches(0.66)
for i, (name, frac, col) in enumerate(teams):
    yb = by0 + i * step
    _, tfl = textbox(s, chart_x + Inches(0.28), yb - Inches(0.02), Inches(1.2), Inches(0.42),
                     anchor=MSO_ANCHOR.MIDDLE)
    add_para(tfl, name, size=10.5, color=SLATE, font=SANS, first=True, space_after=0)
    rect(s, base_x, yb, Emu(int(bar_area * frac)), bar_h, col, radius=True)
footer(s, 10)

# ============================================================ SLIDE 11 — Observability how
s = slide()
kicker(s, "Pillar 03 · Visibility", color=SAGE)
title(s, "Two dashboards: usage you get for free, governance you own")
_, tf = textbox(s, Inches(0.9), Inches(1.9), Inches(11.4), Inches(0.7))
add_para(tf, "The gateway forwards its metrics over native OTLP to CloudWatch, which splits visibility across "
             "two dashboards on purpose — one managed by AWS for usage, one shipped by the stack for governance.",
         size=14, color=SLATE, font=SANS, first=True, space_after=0, line_spacing=1.15)
# two columns
colw = Inches(5.6); c1 = Inches(0.9); c2 = Inches(6.83); cy = Inches(2.75); ch = Inches(3.85)
rect(s, c1, cy, colw, ch, CARD, line=LINE, line_w=Pt(1), radius=True)
_, tf = textbox(s, c1 + Inches(0.32), cy + Inches(0.32), colw - Inches(0.64), ch - Inches(0.6))
add_para(tf, "Coding Agent Insights, managed", size=17, color=SAGE, bold=True, font=SERIF, first=True, space_after=8)
add_para(tf, "Native OTLP metrics auto-populate CloudWatch's built-in Coding Agent Insights dashboard — cost, "
             "tokens, and adoption sliced by user, team, and model. You create nothing; it appears once metrics flow.",
         size=13, color=SLATE, font=SANS, space_after=8, line_spacing=1.18)
add_para(tf, "The same metrics are queryable in PromQL and drive the cost alarms, whether one watches the "
             "organization's total spend or a single individual's.",
         size=13, color=SLATE, font=SANS, space_after=0, line_spacing=1.18)
rect(s, c2, cy, colw, ch, CARD, line=LINE, line_w=Pt(1), radius=True)
_, tf = textbox(s, c2 + Inches(0.32), cy + Inches(0.32), colw - Inches(0.64), ch - Inches(0.6))
add_para(tf, "The governance dashboard, yours", size=17, color=SKY, bold=True, font=SERIF, first=True, space_after=8)
add_para(tf, "The stack ships a lean governance dashboard for the audit signals Coding Agent Insights does not "
             "cover — an event for each tool decision, sign-in, and error, kept as searchable logs.",
         size=13, color=SLATE, font=SANS, space_after=8, line_spacing=1.18)
add_para(tf, "This is where you answer the pointed questions — who was denied a tool, whose sign-in failed, "
             "how spend breaks down by team — by querying the events rather than turning every one into its "
             "own metric.",
         size=13, color=SLATE, font=SANS, space_after=0, line_spacing=1.18)
footer(s, 11)

# ============================================================ SLIDE 12 — defense in depth
s = slide()
kicker(s, "The foundation")
title(s, "Governance rests on identity and isolation")
items = [
    ("Corporate sign-in", "Developers authenticate through your own identity provider, so no API keys or "
        "cloud credentials ever sit on a laptop.", CLAY),
    ("Sessions that expire", "Access is granted as a short-lived token, so removing someone in your "
        "directory ends their gateway access within the session lifetime you set.", SKY),
    ("Private by design", "The gateway is reachable only on your internal network, and Claude Code refuses "
        "to connect to one that resolves to a public address.", SAGE),
    ("Least privilege", "The gateway's own role is allowed to invoke Claude on Bedrock and nothing else, "
        "so a compromise cannot reach the rest of your account.", CLAY_DEEP),
]
cx = Inches(0.9); cw = Inches(2.83); gap = Inches(0.2); cy = Inches(2.3); ch = Inches(3.6)
for i, (h, body, col) in enumerate(items):
    x = cx + i * (cw + gap)
    rect(s, x, cy, cw, ch, CARD, line=LINE, line_w=Pt(1), radius=True)
    rect(s, x + Inches(0.28), cy + Inches(0.32), Inches(0.5), Inches(0.5), col, radius=True)
    _, tf = textbox(s, x + Inches(0.28), cy + Inches(1.05), cw - Inches(0.56), ch - Inches(1.3))
    add_para(tf, h, size=16.5, color=INK, bold=True, font=SERIF, first=True, space_after=8)
    add_para(tf, body, size=12.5, color=SLATE, font=SANS, space_after=0, line_spacing=1.16)
footer(s, 12)

# ============================================================ SLIDE 13 — demo flow
s = slide()
rect(s, 0, 0, Inches(0.28), EMU_H, CLAY)
kicker(s, "Live demo")
title(s, "The five-minute governance walkthrough")
steps = [
    ("1", "Set a deliberately tight cap", "We give a test group a one-dollar daily budget through the admin API."),
    ("2", "Spend against it", "We run a few Claude Code turns and watch the running total climb toward that ceiling."),
    ("3", "Reach the limit", "The next request is refused with a billing error, showing the circuit breaker in action."),
    ("4", "Raise it and recover", "We lift the cap with a single call, and the very next request succeeds without any redeployment."),
    ("5", "Show access and visibility", "We demonstrate a tool withheld from one group, then open the managed Coding Agent Insights dashboard of top spenders."),
]
y = Inches(2.15)
for num, h, body in steps:
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.95), y, Inches(0.62), Inches(0.62))
    circ.fill.solid(); circ.fill.fore_color.rgb = CLAY; circ.line.fill.background(); circ.shadow.inherit = False
    ctf = circ.text_frame; ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run(); cr.text = num; _set_font(cr, size=20, color=IVORY, bold=True, font=SERIF)
    _, tf = textbox(s, Inches(1.85), y - Inches(0.02), Inches(10.4), Inches(0.72), anchor=MSO_ANCHOR.MIDDLE)
    add_para(tf, h, size=17, color=INK, bold=True, font=SANS, first=True, space_after=1)
    add_para(tf, body, size=13, color=SLATE, font=SANS, space_after=0)
    y += Inches(0.92)
_, tf = textbox(s, Inches(1.85), y + Inches(0.0), Inches(10.4), Inches(0.4))
add_para(tf, "The exact commands for each step are in slides/DEMO-RUNBOOK.md.", size=12, color=CLAY,
         bold=True, font=SANS, first=True, space_after=0)
footer(s, 13)

# ============================================================ SLIDE 14 — close
s = slide()
rect(s, 0, 0, Inches(0.28), EMU_H, CLAY)
kicker(s, "In summary")
title(s, "You can configure it, track it, and cap it")
# two-column: what the gateway enforces vs. what it reveals
colw = Inches(5.6); c1 = Inches(0.9); c2 = Inches(6.83); cy = Inches(2.15); ch = Inches(3.6)
rect(s, c1, cy, colw, ch, CARD, line=LINE, line_w=Pt(1), radius=True)
_, tf = textbox(s, c1 + Inches(0.32), cy + Inches(0.34), colw - Inches(0.64), ch - Inches(0.6))
add_para(tf, "What the gateway enforces", size=16, color=SAGE, bold=True, font=SERIF, first=True, space_after=12)
for t in ["It refuses requests once a budget is spent.",
          "It rejects any model a team is not permitted to use.",
          "It withholds the tools you choose to restrict.",
          "It ties every session to your corporate identity."]:
    add_para(tf, "✓   " + t, size=13.5, color=SLATE, font=SANS, space_after=9, line_spacing=1.14)
rect(s, c2, cy, colw, ch, CARD, line=LINE, line_w=Pt(1), radius=True)
_, tf = textbox(s, c2 + Inches(0.32), cy + Inches(0.34), colw - Inches(0.64), ch - Inches(0.6))
add_para(tf, "What the gateway reveals", size=16, color=SKY, bold=True, font=SERIF, first=True, space_after=12)
for t in ["It attributes cost to each person and team.",
          "It presents usage in a managed dashboard, governance in its own.",
          "It keeps a searchable audit trail of what was done.",
          "It alerts you the moment spending looks unusual."]:
    add_para(tf, "✓   " + t, size=13.5, color=SLATE, font=SANS, space_after=9, line_spacing=1.14)
_, tf = textbox(s, Inches(0.9), Inches(6.05), Inches(11.5), Inches(0.7))
add_para(tf, "It is one self-hosted gateway on AWS, running on your credential, under your policy, in full view.",
         size=16, color=INK, italic=True, font=SERIF, first=True, space_after=0, align=PP_ALIGN.CENTER)
footer(s, 14)

# ----------------------------------------------------------------- save
import os
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "claude-gateway-governance.pptx")
prs.save(out)
print(f"wrote {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
